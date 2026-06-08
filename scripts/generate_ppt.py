"""Generate project presentation (10 slides, minimalist style)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "project_presentation.pptx"

# Minimal palette
BG = RGBColor(250, 250, 252)
TITLE = RGBColor(20, 33, 61)
BODY = RGBColor(45, 55, 72)
ACCENT = RGBColor(0, 102, 153)
MUTED = RGBColor(120, 130, 145)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.9))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(12), Inches(0.5))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = MUTED

    line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.05 if subtitle else 1.0), Inches(1.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_bullets(slide, items, top=1.55, left=0.85, width=11.5, height=5.5, font_size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            title, desc = item
            p.text = title
            p.font.bold = True
            p.font.size = Pt(font_size)
            p.font.color.rgb = BODY
            p.space_after = Pt(2)
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(font_size - 1)
            p2.font.color.rgb = MUTED
            p2.level = 1
            p2.space_after = Pt(10)
        else:
            p.text = f"• {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = BODY
            p.space_after = Pt(8)


def slide_title(prs, title, subtitle, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, title, subtitle)
    add_bullets(slide, bullets)


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    t = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.2))
    tp = t.text_frame.paragraphs[0]
    tp.text = "动物图像识别系统"
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = TITLE

    s = slide.shapes.add_textbox(Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.0))
    sp = s.text_frame.paragraphs[0]
    sp.text = "Mask-First 定位 · 多骨干融合 · 不确定性决策"
    sp.font.size = Pt(20)
    sp.font.color.rgb = ACCENT

    m = slide.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.6))
    mp = m.text_frame.paragraphs[0]
    mp.text = "Computer Vision Project · Animal-90 / Multi-Animal Localization"
    mp.font.size = Pt(14)
    mp.font.color.rgb = MUTED

    # Slide 2 — Motivation
    slide_title(
        prs,
        "01  研究动机",
        "Why mask-first instead of whole-image classification?",
        [
            ("问题", "全图分类在复杂背景、多动物同框时，背景像素会干扰特征提取，导致泛化差。"),
            ("Animal-90 局限", "每图仅 1 个动物，定位 IoU 虚高（≈99%），无法检验真实定位能力。"),
            ("目标", "构建「先定位 → 再抑制背景 → 再分类」的两阶段流水线，并在多动物真实/合成数据上训练定位器。"),
            "核心思路：将识别问题分解为 Localization + Classification 两个可独立优化子任务。",
        ],
    )

    # Slide 3 — Pipeline
    slide_title(
        prs,
        "02  整体流程",
        "End-to-end pipeline",
        [
            "输入图像 → 【定位 Localizer】→ 输出 xywh 边界框（归一化坐标）",
            "→ 【Mask】按 bbox 抑制背景 → 【三骨干分类器】→ 【不确定性融合】→ 类别 + 置信度",
            ("LocateAnything（可选）", "推理时优先使用 NVIDIA LocateAnything-3B 零样本定位；不可用时回退自训练 BboxLocalizer。"),
            ("两阶段独立训练", "定位与分类分开反向传播（Separate Backpropagation），各自保存最优 checkpoint，避免任务梯度互相干扰。"),
            "评测指标：Localization IoU · 单模型 val_acc · 融合 ensemble acc(masked/full)",
        ],
    )

    # Slide 4 — Localization
    slide_title(
        prs,
        "03  阶段一：动物定位",
        "BboxLocalizer + Multi-animal training data",
        [
            ("BboxLocalizer", "ResNet18 骨干 + 全连接头，输出 sigmoid 归一化 xywh；损失 = Smooth L1 + GIoU（Generalized IoU，广义交并比）。"),
            ("多动物训练集", "animal-90-multianimal（合成多动物，2700 张）+ Snapshot Serengeti 子集（真实相机陷阱，≥2 动物/bbox 标注）。"),
            ("样本展开", "每张多动物图按 object 展开为 (image, bbox) 对，使定位器学习紧致框而非整图框。"),
            ("评测升级", "除单框 IoU 外，增加 best-of-GT IoU：预测框与任一 GT 框的最大 IoU，反映多动物场景真实表现。"),
        ],
    )

    # Slide 5 — Mask-First
    slide_title(
        prs,
        "04  Mask-First 分类策略",
        "Background suppression before classification",
        [
            ("Mask 操作", "根据定位框生成软/硬矩形 mask，将 bbox 外像素置零：masked = image ⊙ mask。"),
            ("意义", "迫使分类器关注动物区域，减少草地、天空等背景纹理对决策的干扰（Background Suppression）。"),
            ("与全图分类对比", "eval 同时报告 acc(masked) 与 acc(full)，量化 mask 带来的增益。"),
            "分类输入为 224×224、ImageNet 归一化后的 masked crop。",
        ],
    )

    # Slide 6 — ResNet + CBAM
    slide_title(
        prs,
        "05  ResNet50 + CBAM",
        "Convolutional Block Attention Module",
        [
            ("ResNet（残差网络）", "通过 Bottleneck 块与跳跃连接（Skip Connection）缓解深层网络梯度消失；本项目使用 ImageNet 预训练 ResNet50。"),
            ("CBAM 是什么？", "Convolutional Block Attention Module：在特征图上依次施加通道注意力 + 空间注意力，让网络学会「哪些通道重要、哪些位置重要」。"),
            ("Channel Attention（通道注意力）", "对全局平均池化与最大池化结果做 MLP，生成通道权重 — 回答「什么特征（纹理/形状）有用？」"),
            ("Spatial Attention（空间注意力）", "沿通道维做 avg/max，经 7×7 卷积生成空间权重 — 回答「图像哪个区域有用？」"),
            ("与标准 ResNet 的区别", "标准 ResNet 对所有通道/位置等权处理；ResNet+CBAM 在每个 Bottleneck 的 conv3 之后插入 CBAM，自适应重加权特征，更擅长抑制背景噪声。"),
        ],
    )

    # Slide 7 — EfficientNet & ConvNeXt
    slide_title(
        prs,
        "06  EfficientNet-B3 与 ConvNeXt-T",
        "Two complementary backbones vs. ResNet50",
        [
            ("EfficientNet-B3", "采用 Compound Scaling（复合缩放）：同时缩放网络深度、宽度、分辨率；使用 MBConv + SE 模块（Squeeze-and-Excitation，挤压激励通道注意力）。"
             "相比 ResNet：参数/计算更高效，偏「轻量高吞吐」，擅长细粒度纹理。"),
            ("ConvNeXt-T（ConvNeXt-Tiny）", "将 ResNet 式卷积块现代化：大 kernel depthwise conv、LayerNorm、GELU，结构更接近 Vision Transformer 但仍是纯卷积。"
             "相比 ResNet：归纳偏置更现代，全局/context 建模更强，对复杂姿态更鲁棒。"),
            ("为何三模型并行？", "ResNet+CBAM（注意力增强卷积）、EfficientNet（效率+SE）、ConvNeXt（现代卷积）特征空间互补，单模型误差模式不同，适合集成。"),
            "三者均在 ImageNet 预训练权重上微调，仅解冻分类头与末层特征。",
        ],
    )

    # Slide 8 — Ensemble
    slide_title(
        prs,
        "07  不确定性融合集成",
        "Uncertainty-aware fusion (not simple averaging)",
        [
            ("为何不用平均？", "简单 logits/probs 平均假设各模型同等可靠；实际中不同样本、不同模型置信度差异大。"),
            ("Softmax + 熵", "各模型输出经 temperature-scaled softmax 得概率 p_i；预测熵 H = −Σ p log p 衡量不确定性（entropy 越高越「犹豫」）。"),
            ("融合权重", "w_i ∝ exp(−H_i) × max(p_i)：低熵 + 高峰值概率 → 权重更大；再对概率加权求和得到 fused_probs。"),
            ("直觉", "每个样本动态选择「当前最自信」的模型主导决策，而非固定权重投票。"),
            "推理 Pipeline 与 evaluate 脚本均使用 UncertaintyFusionEnsemble 输出最终类别。",
        ],
    )

    # Slide 9 — Training innovations
    slide_title(
        prs,
        "08  训练与评测策略创新",
        "Generalization-oriented design",
        [
            ("随机分层划分", "每次训练用 random split seed 对 Animal-90 做 stratified train/val split；即使 --resume 也重新划分，降低对固定划分的过拟合。"),
            ("多数据源定位", "Localization 不依赖单动物 Animal-90，而合并 multianimal + Serengeti bbox 标注（COCO 格式 xywh）。"),
            ("分阶段 checkpoint", "localizer_best.pth + 三个 classifier_*_best.pth + class_names.json + train_state.json（记录 seed、backbone）。"),
            ("可复现实验", "日志目录记录 download / train / metrics_summary；分别报告单模型 acc 与 ensemble acc。"),
        ],
    )

    # Slide 10 — References
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "09  参考文献", "References")
    refs = [
        "[1] He K. et al. Deep Residual Learning for Image Recognition. CVPR 2016.",
        "[2] Woo S. et al. CBAM: Convolutional Block Attention Module. ECCV 2018.",
        "[3] Tan M., Le Q. EfficientNet: Rethinking Model Scaling for CNNs. ICML 2019.",
        "[4] Liu Z. et al. A ConvNet for the 2020s (ConvNeXt). CVPR 2022.",
        "[5] Hu J. et al. Squeeze-and-Excitation Networks. CVPR 2018.",
        "[6] Swanson A.B. et al. Snapshot Serengeti (camera trap dataset). Sci. Data 2015; Dryad doi:10.5061/dryad.5pt92",
        "[7] LILA BC — Snapshot Serengeti bounding boxes & metadata. https://lila.science/datasets/snapshot-serengeti/",
        "[8] Banerjee S. Animal Image Dataset (90 classes). Kaggle / Hugging Face lucabaggi/animal-wildlife",
        "[9] Rezatofighi H. et al. Generalized Intersection over Union (GIoU). CVPR 2019.",
        "[10] NVIDIA LocateAnything-3B — zero-shot visual grounding (optional inference backend).",
    ]
    add_bullets(slide, refs, top=1.6, font_size=13)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
